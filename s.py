from pptx import Presentation
from pptx.util import Inches, Pt

# Create a PowerPoint presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Title Slide
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "AI-Based Pollution Monitoring & Control System"
slide.placeholders[1].text = "With Integrated Air and Water Purification"

# Helper function to add bullet slides
def add_bullet_slide(title, bullet_points):
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    content.clear()
    for i, point in enumerate(bullet_points):
        p = content.add_paragraph() if i > 0 else content.paragraphs[0]
        p.text = point
        p.level = 0

# Slide content
slides_content = [
    ("Introduction", [
        "Industrialization and urbanization increased pollution levels.",
        "AI and IoT enable real-time monitoring and automated control.",
        "System includes air and water purification actions.",
    ]),
    ("Literature Survey", [
        "AI models predict air quality (Smith et al., 2020).",
        "DL used for source identification (Johnson et al., 2018).",
        "IoT systems detect pollution and trigger filters (Wang & Li, 2017).",
        "DL predicts water quality changes (Kumar et al., 2019).",
        "AI triggers air/water purifiers (Lee & Park, 2021).",
    ]),
    ("Methodology", [
        "IoT sensors collect air, water, and soil data.",
        "AI models (CNN, LSTM) forecast pollution trends.",
        "Real-time alerts activate purification systems.",
        "Air purifiers and water filtration systems reduce pollution levels.",
    ]),
    ("Design and Implementation", [
        "Sensor Layer: IoT sensors (PM, gas, pH, turbidity).",
        "Processing Layer: Cloud AI models analyze trends.",
        "Control Layer: AI triggers purifiers and alerts.",
        "Tech Stack: Arduino, ESP32, AWS, LSTM, RO, UV.",
    ]),
    ("Results and Performance", [
        "Air pollution prediction accuracy: 88%.",
        "Water pollution prediction accuracy: 90%.",
        "Air purification reduced PM10 by 25%.",
        "Water filtration reduced turbidity by 35%.",
        "Alert precision: 85%, recall: 87%.",
    ]),
    ("Future Scope", [
        "Integration with smart cities and edge computing.",
        "Personalized pollution alerts.",
        "Expanded sensor networks.",
        "Renewable-powered purification systems.",
    ]),
    ("Conclusion", [
        "AI enables proactive pollution forecasting and control.",
        "Integrated purification systems mitigate real-time risks.",
        "Hybrid AI models improve system effectiveness.",
        "System is scalable for urban, industrial, and ecological use.",
    ])
]

# Add all slides to the presentation
for title, bullets in slides_content:
    add_bullet_slide(title, bullets)

# Save the presentation
pptx_filename = "AI_Pollution_Monitoring_Purification_System.pptx"
prs.save(pptx_filename)

pptx_filename
